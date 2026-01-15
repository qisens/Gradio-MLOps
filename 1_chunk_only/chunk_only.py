import re
from pptx import Presentation
from langchain_openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.schema import Document
from dotenv import load_dotenv

load_dotenv()

PPT_PATH = "../Gradio MLOps Guide_260101.pptx"

# -----------------------------
# 1. PPT 텍스트 추출
# -----------------------------
def load_ppt(ppt_path):
    prs = Presentation(ppt_path)
    slides = []

    for idx, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)

        slides.append({
            "slide": idx + 1,
            "text": "\n".join(texts)
        })

    return slides


# -----------------------------
# 2. 규칙 기반 청크 생성
# -----------------------------
def chunk_slide(slide):
    documents = []
    lines = [l.strip() for l in slide["text"].split("\n") if l.strip()]

    current_step_text = None
    current_step_num = 0
    buffer = []

    for line in lines:
        # STEP 패턴
        if re.match(r"^(\d+\)|\d+\.)", line):
            # 이전 step 정리
            if current_step_text:
                documents.append(Document(
                    page_content="\n".join([current_step_text] + buffer),
                    metadata={
                        "slide": slide["slide"],
                        "type": "step",
                        "step": current_step_num
                    }
                ))
                buffer = []

            current_step_num += 1
            current_step_text = line

        else:
            # step 이전의 설명은 무시 (섹션 제목 등)
            if current_step_text:
                buffer.append(line)

    # 마지막 step flush
    if current_step_text:
        documents.append(Document(
            page_content="\n".join([current_step_text] + buffer),
            metadata={
                "slide": slide["slide"],
                "type": "step",
                "step": current_step_num
            }
        ))

    return documents



# -----------------------------
# 3. 전체 PPT → 청크
# -----------------------------
def build_chunks(slides):
    all_chunks = []
    for slide in slides:
        slide_chunks = chunk_slide(slide)
        all_chunks.extend(slide_chunks)
    return all_chunks

def debug_print_chunks(documents, limit=50):
    print("\n===== CHUNK STRUCTURE =====\n")

    for i, doc in enumerate(documents[:limit]):
        meta = doc.metadata
        print(f"[{i+1}] "
              f"Slide {meta.get('slide')} | "
              f"Step {meta.get('step')} | "
              f"Type: {meta.get('type')}")
        print(f"  → {doc.page_content}")
        print("-" * 60)

    print(f"\n총 청크 수: {len(documents)}")


# -----------------------------
# 4. 벡터 DB 생성
# -----------------------------
def build_vectorstore(documents):
    embeddings = OpenAIEmbeddings(
        model="text-embedding-3-large"
    )
    return FAISS.from_documents(documents, embeddings)


# -----------------------------
# 5. 검색 테스트
# -----------------------------
def test_search(db):
    print("\n=== 검색 테스트 ===")
    queries = [
        "이미지뷰어 탭에서 추론 결과를 확인하는 방법 순서대로 알려줘.",
        "서버로 설정 버튼 클릭하면 어떻게돼?"
    ]

    for q in queries:
        print(f"\n[Query] {q}")
        docs = db.similarity_search(q, k=3)

        for d in docs:
            print(
                f"- ({d.metadata}) {d.page_content}"
            )


# -----------------------------
# main
# -----------------------------
if __name__ == "__main__":
    slides = load_ppt(PPT_PATH)
    documents = build_chunks(slides)

    debug_print_chunks(documents)
    print(f"총 청크 개수: {len(documents)}")

    db = build_vectorstore(documents)
    test_search(db)

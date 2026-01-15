'''
json을 저장하지 않고 실행때마다 만들어서 사용하기 때문에 시간도 오래걸리고 cost도 높은편
용도에 대한 답은 잘 하지만, 사용법에 대한 답은 잘 하지 못함
'''

import os
import json
import numpy as np
from typing import List, Dict
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation

# =====================
# CLIENT
# =====================
def create_openai_client() -> OpenAI:
    load_dotenv()
    return OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# =====================
# LLM: PPT → MD
# =====================
def extract_text_from_ppt(ppt_path: str) -> str:
    prs = Presentation(ppt_path)
    lines = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

        if not slide_texts:
            continue  # 텍스트 없는 슬라이드는 스킵

        lines.append(f"\n---\n# Slide {i}")
        lines.extend(slide_texts)

    return "\n".join(lines)

# =====================
# LLM: MD → STRUCTURED JSON
# =====================
def convert_md_to_structured_docs(
    client: OpenAI,
    md_text: str
) -> List[Dict]:

    system_prompt = """
너는 기술 문서를 RAG에 적합한 구조화 문서로 변환하는 도우미다.

규칙:
- 각 chunk는 하나의 목적만 가진다
- type은 반드시 다음 중 하나:
  OVERVIEW, SECTION, STEP, RESULT
- STEP: 사용자의 행동
- RESULT: 행동 이후의 화면/상태 변화
- OVERVIEW: 프로그램 전체 설명
- content는 한국어 자연어 문장
- JSON 배열만 출력
"""

    user_prompt = f"""
    아래 markdown 문서를 분석해 구조화하라.

    Markdown:
    {md_text}

    bash
    코드 복사

    반드시 JSON 배열만 출력하라.
    다른 설명, 코드블록, 텍스트는 절대 포함하지 마라.

    출력 예:
    [
      {{
        "type": "OVERVIEW",
        "title": "...",
        "content": "..."
      }}
    ]
    """

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        temperature=0,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
    )

    raw = response.choices[0].message.content
    print("\n[LLM RAW OUTPUT]\n", raw)

    parsed = json.loads(raw)

    # 1️⃣ documents wrapper 제거
    if isinstance(parsed, dict) and "documents" in parsed:
        parsed = parsed["documents"]

    # 2️⃣ 반드시 List로 강제
    if not isinstance(parsed, list):
        raise TypeError(f"Expected list, got {type(parsed)}")

    # 3️⃣ dict만 필터링 + 필수 key 보장
    cleaned_docs = []
    for i, item in enumerate(parsed):
        if not isinstance(item, dict):
            print(f"⚠️ skip non-dict item at index {i}: {item}")
            continue

        if "content" not in item:
            print(f"⚠️ skip item without content at index {i}")
            continue

        cleaned_docs.append({
            "type": item.get("type", "UNKNOWN"),
            "title": item.get("title", ""),
            "content": item["content"]
        })

    if not cleaned_docs:
        raise ValueError("No valid documents after cleaning")

    return cleaned_docs


# =====================
# EMBEDDING
# =====================
def embed_text(client: OpenAI, text: str) -> np.ndarray:
    emb = client.embeddings.create(
        model="text-embedding-3-large",
        input=text
    )
    return np.array(emb.data[0].embedding)


def build_embedding_index(
    client: OpenAI,
    docs: List[Dict]
) -> List[Dict]:

    indexed_docs = []
    for doc in docs:
        vector = embed_text(client, doc["content"])
        indexed_docs.append({**doc, "embedding": vector})

    return indexed_docs


# =====================
# SEARCH
# =====================
def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    return np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b))


def classify_query_type(query: str) -> str:
    overview_keywords = ["뭐", "무엇", "개요", "전체", "설명", "구조"]
    if any(k in query for k in overview_keywords):
        return "OVERVIEW"
    return "DETAIL"


def retrieve_relevant_chunks(
    client: OpenAI,
    query: str,
    docs: List[Dict],
    top_k: int = 5
) -> List[Dict]:

    query_embedding = embed_text(client, query)
    query_type = classify_query_type(query)

    candidates = docs
    if query_type == "OVERVIEW":
        candidates = [d for d in docs if d["type"] == "OVERVIEW"]

    scored = []
    for doc in candidates:
        score = cosine_similarity(query_embedding, doc["embedding"])
        scored.append((score, doc))

    scored.sort(key=lambda x: x[0], reverse=True)
    return [doc for score, doc in scored[:top_k]]


# =====================
# ANSWER GENERATION
# =====================
def generate_answer_from_context(
    client: OpenAI,
    query: str,
    contexts: List[Dict]
) -> str:

    context_block = "\n\n".join(
        f"[{c['type']}] {c['content']}" for c in contexts
    )

    system_prompt = """
너는 사용 설명서 기반 AI 어시스턴트다.
문서에 있는 정보만 사용해서 답변하라.
"""

    user_prompt = f"""
문서 내용:
{context_block}

질문:
{query}

단계가 있으면 단계적으로, 없으면 요약 형태로 답변하라.
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        temperature=0.2,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
    )

    return response.choices[0].message.content


# =====================
# PIPELINE RUNNER
# =====================
def run_rag_pipeline_from_ppt(ppt_path: str):
    client = create_openai_client()
    md_text = extract_text_from_ppt(ppt_path)

    docs = convert_md_to_structured_docs(client, md_text)
    indexed_docs = build_embedding_index(client, docs)

    while True:
        query = input("\n질문 입력 (exit 입력 시 종료): ")
        if query.lower() == "exit":
            break

        retrieved = retrieve_relevant_chunks(client, query, indexed_docs)

        print("\n[검색된 chunk]")
        for r in retrieved:
            print(f"- ({r['type']}) {r['title']}")

        answer = generate_answer_from_context(client, query, retrieved)
        print("\n[답변]")
        print(answer)


# =====================
# ENTRY POINT
# =====================
if __name__ == "__main__":
    ppt_file_path = "./Gradio MLOps Guide_260101.pptx"   # ← 여기서만 설정
    run_rag_pipeline_from_ppt(ppt_file_path)
# ppt => md => json => chunk => embedding
from pptx import Presentation
from pathlib import Path
import re
import os
import json

from openai import OpenAI
from dotenv import load_dotenv
load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def extract_text_from_shape(shape):
    """shape에서 텍스트 추출"""
    if not shape.has_text_frame:
        return None

    lines = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)

    return lines if lines else None


def ppt_to_markdown(ppt_path: str, output_md: str):
    prs = Presentation(ppt_path)
    md_lines = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        md_lines.append(f"\n---\n")
        md_lines.append(f"<!-- Slide {slide_idx} -->\n")

        for shape in slide.shapes:
            text_lines = extract_text_from_shape(shape)
            if not text_lines:
                continue

            # 첫 줄은 제목일 가능성 높음
            if len(text_lines) == 1:
                md_lines.append(f"## {text_lines[0]}\n")
            else:
                for line in text_lines:
                    # 숫자로 시작하면 순서 유지
                    if line.lstrip().startswith(tuple("123456789")):
                        md_lines.append(f"{line}\n")
                    else:
                        md_lines.append(f"- {line}\n")

        md_lines.append("\n")

    Path(output_md).write_text("".join(md_lines), encoding="utf-8")
    print(f"✅ Markdown 파일 생성 완료: {output_md}")

def classify_line(line: str) -> str:
    line = line.strip()
    if not line:
        return None

    # STEP
    if line.endswith("사용법"):
        return "STEP"
    if re.match(r"^\d+\)", line):
        return "STEP"
    if any(k in line for k in ["클릭", "선택", "지정", "버튼", "업로드", "저장"]):
        return "STEP"

    # RESULT
    if line.startswith(("→", "")):
        return "RESULT"
    if any(k in line for k in ["자동", "출력", "불러오기", "그림", "표시"]):
        return "RESULT"

    # FACT
    if any(k in line for k in ["용도", "위함", "기능", "설정", "경로", "Default"]):
        return "FACT"

    return "NOTE"

def structure_md(input_md: str, output_md: str):
    lines = Path(input_md).read_text(encoding="utf-8").splitlines()

    current_title = None
    buffer = {"FACT": [], "STEP": [], "RESULT": [], "NOTE": []}
    out_lines = []

    def flush_section():
        nonlocal buffer
        if not current_title:
            return

        out_lines.append(f"\n## {current_title}\n")

        for key in ["FACT", "STEP", "RESULT", "NOTE"]:
            if buffer[key]:
                out_lines.append(f"[{key}]\n")
                for item in buffer[key]:
                    out_lines.append(f"- {item}\n")
                out_lines.append("\n")

        buffer = {"FACT": [], "STEP": [], "RESULT": [], "NOTE": []}

    for line in lines:
        # 슬라이드 구분
        if line.startswith("---"):
            flush_section()
            continue

        # 제목
        if line.startswith("## "):
            flush_section()
            current_title = line.replace("## ", "").strip()
            continue

        # 주석 무시
        if line.startswith("<!--"):
            continue

        tag = classify_line(line)
        if tag:
            buffer[tag].append(line.lstrip("- ").strip())

    flush_section()
    Path(output_md).write_text("".join(out_lines), encoding="utf-8")
    print(f"✅ 구조화 md 생성 완료: {output_md}")

def parse_structured_md(md_text: str):
    chunks = []
    current_title = None
    current_type = None
    buffer = []

    for line in md_text.splitlines():
        line = line.strip()

        if line.startswith("## "):
            current_title = line.replace("## ", "").strip()
            continue

        if line.startswith("[") and line.endswith("]"):
            if buffer:
                chunks.append({
                    "title": current_title,
                    "type": current_type,
                    "content": " ".join(buffer)
                })
                buffer = []
            current_type = line.strip("[]")
            continue

        if line.startswith("- "):
            buffer.append(line.replace("- ", ""))

    if buffer:
        chunks.append({
            "title": current_title,
            "type": current_type,
            "content": " ".join(buffer)
        })

    return chunks


def md_to_json(INPUT_MD, OUTPUT_JSON):
    md_text = Path(INPUT_MD).read_text(encoding="utf-8")
    chunks = parse_structured_md(md_text)

    Path(OUTPUT_JSON).write_text(
        json.dumps(chunks, ensure_ascii=False, indent=2),
        encoding="utf-8"
    )

    print(f"✅ chunk 생성 완료: {OUTPUT_JSON}")
    print(f"총 {len(chunks)}개 chunk")

def embed_text(text: str):
    response = client.embeddings.create(
        model="text-embedding-3-large",
        input=text
    )
    return response.data[0].embedding


def embedding_chunk(INPUT_JSON, OUTPUT_JSON):
    chunks = json.loads(Path(INPUT_JSON).read_text(encoding="utf-8"))

    for chunk in chunks:
        text = f"{chunk['title']} [{chunk['type']}] {chunk['content']}"
        chunk["embedding"] = embed_text(text)

    Path(OUTPUT_JSON).write_text(
        json.dumps(chunks, ensure_ascii=False),
        encoding="utf-8"
    )

    print(f"✅ embedding 완료: {OUTPUT_JSON}")

if __name__ == "__main__":
    ''' step 1 '''
    # ppt를 md 파일로 만드는 과정에서 확인할것
    # 1. 제목에 너무 ###이 많은지
    # 2. bullet이 문장 단위로 잘 잘렸는지
    # 3. 슬라이드 구분 ---- 이 잘 되었는지
    PPT_PATH = "./Gradio MLOps Guide_260101.pptx"
    OUTPUT_MD_step1 = "guide_step1_ppt_to_md.md"
    if not os.path.exists(OUTPUT_MD_step1):
        ppt_to_markdown(PPT_PATH, OUTPUT_MD_step1)

    ''' step 2 '''
    # step1에서 나온 md 파일을 사람이 직접 검토
    INPUT_MD_step2 = "guide_step2_human.md"
    if not os.path.exists(INPUT_MD_step2):
        print("사람이 직접 step1 문서를 검토하여 md 파일을 수정해 주세요")

    else:
        ''' step 3 '''
        # md를 의미구조 md로 정리하여 fact / step 등으로 구분 : 이게 청크가 될거
        OUTPUT_MD_step3 = "guide_step3_structure.md"
        structure_md(INPUT_MD_step2, OUTPUT_MD_step3)

        ''' step 4 '''
        # 청크 단위로 json을 생성 (이 문서로 임베딩 하여 RAG 검색에 참조함)
        OUTPUT_JSON_step4 = "guide_step4_md_to_chunk.json"
        md_to_json(OUTPUT_MD_step3, OUTPUT_JSON_step4)

        ''' step 5 '''
        # 청크를 임베딩하여 임베딩 json 파일을 만듬
        # md에서 만든 json을 임베딩
        OUTPUT_JSON_step5 = "guide_step5_embedded_chunk.json"
        if not os.path.exists(OUTPUT_JSON_step5):
            embedding_chunk(OUTPUT_JSON_step4, OUTPUT_JSON_step5)
